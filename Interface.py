import tkinter as tk
from tkinter import ttk, filedialog, messagebox, simpledialog, colorchooser
from tkinter import font as tkfont
from PIL import Image, ImageTk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import json
import os
from collections import OrderedDict

# ==== Configuration ====

APP_TITLE = "Compétences Pro Ultimate"
MAX_LINES_PER_SLIDE = 20    # lignes (en-têtes + compétences) par diapositive (utilisé pour l'aperçu)
LEFT_PANEL_MINW = 300       # min largeur colonne gauche
COMP_LISTBOX_WIDTH = 62     # largeur listbox (caractères)
PREVIEW_WIDTH = 900         # valeurs initiales (taille réelle prise au runtime)
PREVIEW_HEIGHT = 520
HEADER_HEIGHT = 48
SUBHEADER_SPACING = 8
LINE_SPACING = 6
COVER_HEADER_HEIGHT = 64
TEXT_MARGIN_X = 24
TEXT_MARGIN_Y = 18
DEFAULT_BODY_FONT = ("Arial", 12)
DEFAULT_TITLE_SIZE_PT = 20
DEFAULT_BODY_SIZE_PT = 12
DEFAULT_SUBHEADER_BOLD = True
DEFAULT_SUBHEADER_UNDERLINE = True
DEFAULT_TITLE_FG = "white"
COVER_HEADER_COLOR = "#6e6e6e"    # gris bandeau couverture
COVER_PERSONAL_BG_PREVIEW = "#6B8E23"

# Palette de couleurs pour domaines (assignation automatique)
DOMAIN_COLORS = [
    "#2E86C1", "#AF7AC5", "#48C9B0", "#F5B041", "#EC7063",
    "#16A085", "#5D6D7E", "#CA6F1E", "#7D3C98", "#1F618D"
]

# Sections (maternelle)
SECTION_KEYS = ["TPS", "PS", "MS", "GS"]
SECTION_LABELS = {
    "TPS": "TOUTE PETITE SECTION",
    "PS": "PETITE SECTION",
    "MS": "MOYENNE SECTION",
    "GS": "GRANDE SECTION",
}
# Champs par section
SECTION_FIELDS = OrderedDict([
    ("annee", "Année scolaire"),
    ("ecole", "École"),
    ("enseignants", "Enseignant(s)"),
])


# ==== Structures de données ====

class DomainState:
    def __init__(self, name, color):
        self.name = name
        self.color = color  # hex
        self.font_body = DEFAULT_BODY_FONT


class CompetenceItem:
    def __init__(self, domain, subdomain, text, ts=None, batch_id=None):
        self.domain = domain
        self.subdomain = subdomain
        self.text = text
        # Ajouts: horodatage & lot d'ajout (pour regrouper dans le PPT)
        self.ts = ts    # "Mois Année" ex: "Février 2023"
        self.batch_id = batch_id  # entier, incrémenté à chaque clic "Ajouter ->"

    def key(self):
        # clé d'unicité
        return (self.domain, self.subdomain or "", self.text)


# ==== Application ====

class CompetenceApp:
    def __init__(self, root):
        self.root = root
        self.root.title(APP_TITLE)

        # Etat
        self.available = OrderedDict()    # domain -> OrderedDict{subdomain -> [competences]}
        self.domain_order = []
        self.domain_states = {}    # domain -> DomainState
        self.selected_items = []    # list[CompetenceItem]
        self.added_set = set()    # keys pour anti-doublon
        self.add_batch_counter = 0

        # Aperçu global
        self.domain_page_map = {}    # domain -> list[page]
        self.item_page_index = {}    # item.key() -> (domain, page_index)
        self.flat_pages = []    # list of (domain, page_index)
        self.current_flat_index = 0
        self.current_domain = None

        # Images par page
        self.page_images = {}    # (domain, page_index) -> [img dict]

        # Infos couverture
        self.nom_var = tk.StringVar()
        self.prenom_var = tk.StringVar()
        self.naissance_var = tk.StringVar()
        self.photo_path = None
        self.personal_completed = False

        # Horodateur (obligatoire)
        self.month_var = tk.StringVar()
        self.year_var = tk.StringVar()

        # Sections (TPS/PS/MS/GS)
        self.sections_data = {
            key: {
                "completed": False,
                "fields": {fname: "" for fname in SECTION_FIELDS.keys()},
                "photo": None,   # chemin photo de la section
                "bilan1": "",
                "bilan2": "",
                "bilan2_enabled": False,
            } for key in SECTION_KEYS
        }
        self.sections_widgets = {}

        # Descriptions domaines/sous-domaines (DOMAINES.txt)
        self.domain_descriptions = {}    # domain -> str
        self.subdomain_descriptions = {}  # (domain, subdomain) -> str

        # Pour mesure du texte
        self.measure_font = tkfont.Font(family="Arial", size=12)

        # Drag/Resize images (aperçu)
        self.drag_data = {"x": 0, "y": 0, "image_index": None}
        self.resize_data = {"image_index": None, "start_x": 0, "start_y": 0}

        # UI
        self._build_ui()
        self.update_cover_preview()
        self.rebuild_pages_and_refresh()

    # ---- UI ----

    def _build_ui(self):
        self.root.geometry("1600x1000")
        self.root.minsize(1200, 800)

        # Conteneur scrollable principal (scrollbar verticale à droite)
        outer = ttk.Frame(self.root)
        outer.pack(fill="both", expand=True)

        self.main_canvas = tk.Canvas(outer, highlightthickness=0)
        vscroll = ttk.Scrollbar(outer, orient="vertical", command=self.main_canvas.yview)
        self.main_canvas.configure(yscrollcommand=vscroll.set)

        vscroll.pack(side="right", fill="y")
        self.main_canvas.pack(side="left", fill="both", expand=True)

        # Frame de contenu interne au canvas
        self.content = ttk.Frame(self.main_canvas)
        self._content_window = self.main_canvas.create_window((0, 0), window=self.content, anchor="nw")

        # Ajuste la scrollregion quand le contenu change
        def _update_scrollregion(event=None):
            self.main_canvas.configure(scrollregion=self.main_canvas.bbox("all"))

        self.content.bind("<Configure>", _update_scrollregion)

        # Assure que la largeur du contenu suit la largeur du canvas (pas de scroll horizontal)
        def _sync_content_width(event):
            self.main_canvas.itemconfigure(self._content_window, width=event.width)

        self.main_canvas.bind("<Configure>", _sync_content_width)

        # Ligne haute: informations personnelles + Sections onglets
        top = ttk.Frame(self.content)
        top.pack(fill="x", padx=8, pady=6)

        # Informations personnelles + Horodatage
        pers = ttk.LabelFrame(top, text="Informations personnelles & Horodatage (obligatoire)")
        pers.pack(side="left", padx=6, pady=4, fill="x", expand=True)

        ttk.Label(pers, text="Nom:").grid(row=0, column=0, sticky="w")
        ttk.Entry(pers, textvariable=self.nom_var, width=18).grid(row=0, column=1, sticky="we", padx=4)

        ttk.Label(pers, text="Prénom:").grid(row=0, column=2, sticky="w")
        ttk.Entry(pers, textvariable=self.prenom_var, width=18).grid(row=0, column=3, sticky="we", padx=4)

        ttk.Label(pers, text="Date de naissance:").grid(row=0, column=4, sticky="w")
        ttk.Entry(pers, textvariable=self.naissance_var, width=16).grid(row=0, column=5, sticky="we", padx=4)

        ttk.Button(pers, text="Importer photo", command=self._import_photo).grid(row=0, column=6, padx=6)
        ttk.Button(pers, text="Marquer comme complétée", command=self._mark_personal_completed).grid(row=0, column=7, padx=6)

        # Horodatage
        ttk.Label(pers, text="Mois (ex: Février):").grid(row=1, column=0, sticky="w", pady=(6, 0))
        ttk.Entry(pers, textvariable=self.month_var, width=14).grid(row=1, column=1, sticky="we", padx=4, pady=(6, 0))
        ttk.Label(pers, text="Année (ex: 2023):").grid(row=1, column=2, sticky="w", pady=(6, 0))
        ttk.Entry(pers, textvariable=self.year_var, width=10).grid(row=1, column=3, sticky="we", padx=4, pady=(6, 0))

        for col in range(8):
            pers.grid_columnconfigure(col, weight=1)

        # Bloc Sections (TPS/PS/MS/GS)
        sections_block = ttk.LabelFrame(self.content, text="Sections de cycle (mémorisées indépendamment)")
        sections_block.pack(fill="x", padx=8, pady=4)

        self.sections_nb = ttk.Notebook(sections_block)
        self.sections_nb.pack(fill="x", padx=6, pady=6)
        for key in SECTION_KEYS:
            self._build_section_tab(key)

        # Zone principale: 3 colonnes proportionnelles
        main = ttk.Frame(self.content)
        main.pack(fill="both", expand=True, padx=8, pady=(4, 0))

        cols = ttk.Panedwindow(main, orient=tk.HORIZONTAL)
        cols.pack(fill="both", expand=True)

        # 1) Colonne gauche: Disponibles
        col_left = ttk.LabelFrame(cols, text="Compétences disponibles")
        cols.add(col_left, weight=1)

        left_tree_frame = ttk.Frame(col_left)
        left_tree_frame.pack(fill="both", expand=True, padx=6, pady=6)

        self.tree = ttk.Treeview(left_tree_frame, show="tree")
        self.tree.pack(side="left", fill="both", expand=True)
        self.tree.column("#0", width=420, minwidth=LEFT_PANEL_MINW, stretch=True)

        yscroll_tree = ttk.Scrollbar(left_tree_frame, orient="vertical", command=self.tree.yview)
        yscroll_tree.pack(side="left", fill="y")
        self.tree.configure(yscrollcommand=yscroll_tree.set)

        # 2) Colonne milieu: Compétences du sous-domaine
        col_mid = ttk.LabelFrame(cols, text="Compétences du sous-domaine")
        cols.add(col_mid, weight=1)

        mid_inner = ttk.Frame(col_mid)
        mid_inner.pack(fill="both", expand=True, padx=6, pady=6)

        self.comps_list = tk.Listbox(mid_inner, width=COMP_LISTBOX_WIDTH, selectmode=tk.EXTENDED)
        self.comps_list.pack(side="left", fill="both", expand=True)

        yscroll_comp = ttk.Scrollbar(mid_inner, orient="vertical", command=self.comps_list.yview)
        yscroll_comp.pack(side="left", fill="y")
        self.comps_list.configure(yscrollcommand=yscroll_comp.set)

        # Boutons bas de colonne milieu
        btns_left = ttk.Frame(col_mid)
        btns_left.pack(fill="x", padx=6, pady=(0, 8))
        ttk.Button(btns_left, text="Charger COMPETENCES.txt", command=self.load_competences_file).pack(side="left", padx=2)
        ttk.Button(btns_left, text="Ajouter ->", command=self.add_selected_competences).pack(side="left", padx=2)

        # 3) Colonne droite: Sélectionnées (PPT)
        col_right = ttk.LabelFrame(cols, text="Sélectionnées (dans le PPT)")
        cols.add(col_right, weight=1)

        sel_area = ttk.Frame(col_right)
        sel_area.pack(fill="both", expand=True, padx=6, pady=6)

        self.selected_tree = ttk.Treeview(sel_area, columns=("subdomain", "text"), show="headings")
        self.selected_tree.heading("subdomain", text="Sous-domaine")
        self.selected_tree.heading("text", text="Compétence")
        self.selected_tree.column("subdomain", width=200, stretch=True)
        self.selected_tree.column("text", width=460, stretch=True)
        self.selected_tree.pack(side="left", fill="both", expand=True)

        yscroll_sel = ttk.Scrollbar(sel_area, orient="vertical", command=self.selected_tree.yview)
        yscroll_sel.pack(side="left", fill="y")
        self.selected_tree.configure(yscrollcommand=yscroll_sel.set)

        # Barre de boutons en bas de la colonne droite (visibles et regroupés sous la liste)
        btns_center = ttk.Frame(col_right)
        btns_center.pack(fill="x", padx=6, pady=(0, 8))
        # Rangée 1
        row1 = ttk.Frame(btns_center)
        row1.pack(fill="x", pady=2)
        ttk.Button(row1, text="Retirer <-", command=self.remove_selected_from_ppt).pack(side="left", padx=2)
        ttk.Button(row1, text="Aller à la page", command=self.goto_selected_page).pack(side="left", padx=2)
        ttk.Button(row1, text="Exporter PowerPoint", command=self.export_ppt).pack(side="right", padx=2)

        # Rangée 2 (autres actions)
        row2 = ttk.Frame(btns_center)
        row2.pack(fill="x", pady=2)
        ttk.Button(row2, text="Ajouter image (page)", command=self.add_image_page).pack(side="left", padx=2)
        ttk.Button(row2, text="Police/Couleur (domaine)", command=self.change_font_color).pack(side="left", padx=2)
        ttk.Button(row2, text="Sauvegarder projet", command=self.save_project).pack(side="left", padx=2)
        ttk.Button(row2, text="Charger projet", command=self.load_project).pack(side="left", padx=2)

        # En dessous des 3 colonnes: zone de prévisualisation globale
        bottom = ttk.Frame(self.content)
        bottom.pack(fill="both", expand=True, padx=8, pady=8)

        cover_frame = ttk.LabelFrame(bottom, text="Mini-aperçu Page de garde")
        cover_frame.pack(fill="x")
        # Hauteur augmentée + redraw auto
        self.cover_canvas = tk.Canvas(cover_frame, height=260, bg="white", highlightthickness=1, highlightbackground="#ddd")
        self.cover_canvas.pack(fill="x")
        self.cover_canvas.bind("<Configure>", lambda e: self.update_cover_preview())

        pager = ttk.Frame(bottom)
        pager.pack(fill="x", pady=6)
        self.btn_prev = ttk.Button(pager, text="◀ Page précédente", command=self.prev_page)
        self.btn_prev.pack(side="left", padx=4)
        self.page_var = tk.StringVar(value="Page 0/0")
        ttk.Label(pager, textvariable=self.page_var).pack(side="left", padx=10)
        self.btn_next = ttk.Button(pager, text="Page suivante ▶", command=self.next_page)
        self.btn_next.pack(side="left", padx=4)

        self.preview_frame = ttk.LabelFrame(bottom, text="Aperçu des pages (tous domaines)")
        self.preview_frame.pack(fill="both", expand=True)
        self.preview_canvas = tk.Canvas(
            self.preview_frame,
            width=PREVIEW_WIDTH, height=PREVIEW_HEIGHT,
            bg="white", highlightthickness=1, highlightbackground="#ddd"
        )
        self.preview_canvas.pack(fill="both", expand=True)
        self.preview_canvas.bind("<Configure>", lambda e: self.update_preview())

        # drag/resize images sur page courante
        self.preview_canvas.bind("<Button-1>", self.start_drag)
        self.preview_canvas.bind("<B1-Motion>", self.drag_image)
        self.preview_canvas.bind("<Button-3>", self.start_resize)
        self.preview_canvas.bind("<B3-Motion>", self.resize_image)

        # Events
        self.tree.bind("<<TreeviewSelect>>", self.on_tree_select)

    def _build_section_tab(self, key):
        title = SECTION_LABELS[key]
        frame = ttk.Frame(self.sections_nb)
        self.sections_nb.add(frame, text=title)

        # Case à cocher "Section complétée"
        completed_var = tk.BooleanVar(value=False)
        chk = ttk.Checkbutton(
            frame,
            text="Marquer cette section comme complétée",
            variable=completed_var,
            command=lambda k=key: self._toggle_section_completed(k)
        )
        chk.grid(row=0, column=0, columnspan=4, sticky="w", pady=(6, 2))

        entries = {}
        r = 1
        for fname, flabel in SECTION_FIELDS.items():
            ttk.Label(frame, text=flabel + " :").grid(row=r, column=0, sticky="e", padx=(0, 6), pady=(6, 2))
            var = tk.StringVar()
            ent = ttk.Entry(frame, width=32, textvariable=var)
            ent.grid(row=r, column=1, sticky="we", pady=(6, 2))

            def on_change(var=var, k=key, fn=fname):
                self.sections_data[k]["fields"][fn] = var.get().strip()
                # auto: tous les champs remplis -> completed True
                all_filled = all(self.sections_data[k]["fields"][f].strip() for f in SECTION_FIELDS.keys())
                self.sections_widgets[k]["completed_var"].set(all_filled)
                self.sections_data[k]["completed"] = all_filled
                self.update_cover_preview()

            var.trace_add("write", lambda *args, cb=on_change: cb())
            entries[fname] = (ent, var)
            r += 1

        # Boutons et photo de la section
        ttk.Button(frame, text="Importer photo de la section", command=lambda k=key: self._import_section_photo(k)).grid(row=r, column=0, sticky="e", pady=6)
        photo_label = ttk.Label(frame, text="Aucune photo")
        photo_label.grid(row=r, column=1, sticky="w", pady=6)
        r += 1

        btns = ttk.Frame(frame)
        btns.grid(row=r, column=0, columnspan=4, sticky="w", pady=8)

        ttk.Button(btns, text="Marquer comme complétée",
                   command=lambda k=key: self._set_section_completed(k, True)).pack(side="left", padx=2)
        ttk.Button(btns, text="Décocher",
                   command=lambda k=key: self._set_section_completed(k, False)).pack(side="left", padx=2)
        ttk.Button(btns, text="Effacer le contenu",
                   command=lambda k=key: self._clear_section(k)).pack(side="left", padx=8)

        # Bilans
        ttk.Button(btns, text="Ajouter bilan",
                   command=lambda k=key: self._add_bilan(k, which=1)).pack(side="left", padx=8)

        bilan2_var = tk.BooleanVar(value=False)
        chk2 = ttk.Checkbutton(btns, text="2e bilan", variable=bilan2_var,
                               command=lambda k=key: self._toggle_bilan2(k))
        chk2.pack(side="left", padx=8)

        bilan2_btn = ttk.Button(btns, text="Ajouter bilan 2",
                                command=lambda k=key: self._add_bilan(k, which=2))
        bilan2_btn.state(["disabled"])
        bilan2_btn.pack(side="left", padx=8)

        frame.grid_columnconfigure(1, weight=1)
        self.sections_widgets[key] = {
            "completed_var": completed_var,
            "entries": entries,
            "photo_label": photo_label,
            "bilan2_var": bilan2_var,
            "bilan2_btn": bilan2_btn,
        }

    # ---- Chargement / parsing ----

    def load_competences_file(self):
        path = filedialog.askopenfilename(
            title="Sélectionnez le fichier COMPETENCES.txt",
            filetypes=[("Fichiers texte", "*.txt")]
        )
        if not path:
            return

        self.available.clear()
        self.domain_order.clear()
        self.domain_states.clear()
        self.selected_items.clear()
        self.added_set.clear()
        self.domain_page_map.clear()
        self.item_page_index.clear()
        self.flat_pages.clear()
        self.page_images.clear()
        self.current_flat_index = 0
        self.current_domain = None

        current_domain = None
        current_subdomain = None

        with open(path, "r", encoding="utf-8-sig") as f:
            for raw in f:
                line = raw.strip()
                if not line:
                    continue
                # Nettoyage unicode
                line = line.replace("\u202f", " ").replace("\u00a0", " ").replace("\u2019", "'")

                if line.startswith("##-Domaine"):
                    current_domain = line.replace("##-Domaine", "").strip()
                    if current_domain not in self.available:
                        self.available[current_domain] = OrderedDict()
                        self.domain_order.append(current_domain)
                        color = DOMAIN_COLORS[(len(self.domain_order)-1) % len(DOMAIN_COLORS)]
                        self.domain_states[current_domain] = DomainState(current_domain, color)
                    current_subdomain = None

                elif line.startswith("#-"):
                    sub = line.replace("#-", "").strip()
                    if sub.lower().startswith("sous-domaine:"):
                        sub = sub.split(":", 1)[1].strip()
                    current_subdomain = sub

                elif line.startswith("XX"):
                    comp = line.replace("XX", "").strip()
                    if current_domain is None:
                        current_domain = "Domaine"
                        if current_domain not in self.available:
                            self.available[current_domain] = OrderedDict()
                            self.domain_order.append(current_domain)
                            color = DOMAIN_COLORS[(len(self.domain_order)-1) % len(DOMAIN_COLORS)]
                            self.domain_states[current_domain] = DomainState(current_domain, color)
                    sd = current_subdomain if current_subdomain else current_domain
                    self.available[current_domain].setdefault(sd, [])
                    self.available[current_domain][sd].append(comp)

        self.build_available_tree()
        self.rebuild_pages_and_refresh()

    def build_available_tree(self):
        self.tree.delete(*self.tree.get_children())
        for domain in self.domain_order:
            d_id = self.tree.insert("", "end", text=domain, open=True)
            submap = self.available[domain]
            if not submap:
                self.tree.insert(d_id, "end", text=domain)
            else:
                for sub in submap.keys():
                    self.tree.insert(d_id, "end", text=sub)

    def on_tree_select(self, event):
        # Affiche les compétences du sous-domaine choisi (retire celles déjà ajoutées)
        self.comps_list.delete(0, tk.END)
        sel = self.tree.selection()
        if not sel:
            return
        item_id = sel[0]
        text = self.tree.item(item_id, "text")
        parent = self.tree.parent(item_id)

        if parent:
            domain = self.tree.item(parent, "text")
            sub = text
            comps = self.available.get(domain, {}).get(sub, [])
            for c in comps:
                key = (domain, sub or "", c)
                if key not in self.added_set:
                    self.comps_list.insert(tk.END, c)
                else:
                    pass

    # ---- Ajout / retrait ----

    def add_selected_competences(self):
        # Vérifier horodatage obligatoire
        month = self.month_var.get().strip()
        year = self.year_var.get().strip()
        if not month or not year:
            messagebox.showinfo("Champs requis", "Veuillez remplir les champs Mois et Année avant d'ajouter des compétences.")
            return
        ts = f"{month} {year}"

        sel = self.tree.selection()
        if not sel:
            messagebox.showinfo("Info", "Sélectionnez un sous-domaine.")
            return
        item_id = sel[0]
        parent = self.tree.parent(item_id)
        if not parent:
            messagebox.showinfo("Info", "Sélectionnez un sous-domaine, pas un domaine.")
            return

        domain = self.tree.item(parent, "text")
        sub = self.tree.item(item_id, "text")
        chosen = [self.comps_list.get(i) for i in self.comps_list.curselection()]
        if not chosen:
            messagebox.showinfo("Info", "Sélectionnez au moins une compétence.")
            return

        # Nouveau lot d'ajout
        self.add_batch_counter += 1
        batch_id = self.add_batch_counter

        added_any = False
        for comp in chosen:
            item = CompetenceItem(domain, sub, comp, ts=ts, batch_id=batch_id)
            if item.key() in self.added_set:
                continue
            self.selected_items.append(item)
            self.added_set.add(item.key())
            added_any = True

        if not added_any:
            return

        self.refresh_selected_tree()
        self.rebuild_pages_and_refresh()
        self.on_tree_select(None)

    def refresh_selected_tree(self):
        self.selected_tree.delete(*self.selected_tree.get_children())
        for it in self.selected_items:
            self.selected_tree.insert("", "end", values=(it.subdomain, it.text))

    def remove_selected_from_ppt(self):
        sel = self.selected_tree.selection()
        if not sel:
            return
        indices = []
        for iid in sel:
            sd, txt = self.selected_tree.item(iid, "values")
            for idx, it in enumerate(self.selected_items):
                if it.subdomain == sd and it.text == txt:
                    indices.append(idx)
                    break
        if not indices:
            return

        for idx in sorted(indices, reverse=True):
            k = self.selected_items[idx].key()
            if k in self.added_set:
                self.added_set.remove(k)
            self.selected_items.pop(idx)

        self.refresh_selected_tree()
        self.rebuild_pages_and_refresh()
        self.on_tree_select(None)

    def goto_selected_page(self):
        sel = self.selected_tree.selection()
        if not sel:
            return
        iid = sel[0]
        sd, txt = self.selected_tree.item(iid, "values")
        for it in self.selected_items:
            if it.subdomain == sd and it.text == txt:
                loc = self.item_page_index.get(it.key())
                if loc:
                    domain, page_idx = loc
                    for i, (d, p) in enumerate(self.flat_pages):
                        if d == domain and p == page_idx:
                            self.current_flat_index = i
                            self.current_domain = d
                            break
                    self.update_preview()
                    return

    # ---- Images (par page) ----

    def _current_page_key(self):
        if not self.flat_pages:
            return None
        return self.flat_pages[self.current_flat_index]  # (domain, page_index)

    def add_image_page(self):
        # Ajoute des images à la page actuellement visible dans l'aperçu
        key = self._current_page_key()
        if not key:
            messagebox.showinfo("Info", "Aucune page n'est disponible.")
            return
        domain, page_index = key
        paths = filedialog.askopenfilenames(
            title="Sélectionnez une ou plusieurs images",
            filetypes=[("Images", "*.png;*.jpg;*.jpeg;*.bmp")]
        )
        if not paths:
            return
        imgs = self.page_images.setdefault((domain, page_index), [])
        for p in paths:
            try:
                pil = Image.open(p)
                pil.thumbnail((360, 360), Image.LANCZOS)
                tkimg = ImageTk.PhotoImage(pil)
                imgs.append({
                    "path": p,
                    "pil": pil,
                    "tk": tkimg,
                    "pos": [60, HEADER_HEIGHT + 30],
                    "size": list(pil.size)
                })
            except Exception as e:
                messagebox.showerror("Image", f"Erreur avec {p}: {e}")
        self.update_preview()

    def change_font_color(self):
        if not self.flat_pages:
            return
        domain, _ = self.flat_pages[self.current_flat_index]
        ds = self.domain_states.get(domain)
        if not ds:
            return
        size = simpledialog.askinteger("Taille police corps", "Taille de la police (ex: 12):",
                                       initialvalue=ds.font_body[1])
        color = colorchooser.askcolor(title="Couleur du domaine (bandeau/titres)")[1]
        if size:
            ds.font_body = (ds.font_body[0], size)
        if color:
            ds.color = color
        self.update_preview()

    # Drag & drop / resize (aperçu)
    def _hit_test_image(self, event):
        key = self._current_page_key()
        if not key:
            return None, None
        imgs = self.page_images.get(key, [])
        for i, img in enumerate(reversed(imgs)):
            # sélectionner l'image au-dessus si superposition
            real_index = len(imgs) - 1 - i
            x, y = img["pos"]
            w, h = img["size"]
            if x <= event.x <= x + w and y <= event.y <= y + h:
                return imgs, real_index
        return imgs, None

    def start_drag(self, event):
        imgs, idx = self._hit_test_image(event)
        if idx is None:
            return
        self.drag_data = {"image_index": idx, "x": event.x, "y": event.y}

    def drag_image(self, event):
        key = self._current_page_key()
        if not key:
            return
        idx = self.drag_data.get("image_index")
        if idx is None:
            return
        imgs = self.page_images.get(key, [])
        if idx >= len(imgs):
            return
        img = imgs[idx]
        dx = event.x - self.drag_data["x"]
        dy = event.y - self.drag_data["y"]
        img["pos"][0] += dx
        img["pos"][1] += dy
        self.drag_data["x"] = event.x
        self.drag_data["y"] = event.y
        self.update_preview()

    def start_resize(self, event):
        imgs, idx = self._hit_test_image(event)
        if idx is None:
            return
        self.resize_data = {"image_index": idx, "start_x": event.x, "start_y": event.y}

    def resize_image(self, event):
        key = self._current_page_key()
        if not key:
            return
        idx = self.resize_data.get("image_index")
        if idx is None:
            return
        imgs = self.page_images.get(key, [])
        if idx >= len(imgs):
            return
        img = imgs[idx]
        dx = event.x - self.resize_data["start_x"]
        dy = event.y - self.resize_data["start_y"]
        new_w = max(30, img["size"][0] + dx)
        new_h = max(30, img["size"][1] + dy)
        try:
            pil = Image.open(img["path"]).resize((int(new_w), int(new_h)), Image.LANCZOS)
            img["pil"] = pil
            img["tk"] = ImageTk.PhotoImage(pil)
            img["size"] = [int(new_w), int(new_h)]
            self.resize_data["start_x"] = event.x
            self.resize_data["start_y"] = event.y
            self.update_preview()
        except Exception as e:
            messagebox.showerror("Redimensionner", str(e))

    # ---- Sections - logique ----

    def _toggle_section_completed(self, key):
        val = self.sections_widgets[key]["completed_var"].get()
        self.sections_data[key]["completed"] = bool(val)
        self.update_cover_preview()

    def _set_section_completed(self, key, value):
        self.sections_widgets[key]["completed_var"].set(bool(value))
        self.sections_data[key]["completed"] = bool(value)
        self.update_cover_preview()

    def _clear_section(self, key):
        for fname in SECTION_FIELDS.keys():
            ent, var = self.sections_widgets[key]["entries"][fname]
            var.set("")
        # reset photo
        self.sections_data[key]["photo"] = None
        self.sections_widgets[key]["photo_label"].configure(text="Aucune photo")
        # bilans
        self.sections_data[key]["bilan1"] = ""
        self.sections_data[key]["bilan2"] = ""
        self.sections_data[key]["bilan2_enabled"] = False
        self.sections_widgets[key]["bilan2_var"].set(False)
        self.sections_widgets[key]["bilan2_btn"].state(["disabled"])

        # recalcul auto completed
        self._recalc_section_completed(key)
        self.update_cover_preview()

    def _recalc_section_completed(self, key):
        all_filled = all(self.sections_data[key]["fields"][f].strip() for f in SECTION_FIELDS.keys())
        self.sections_widgets[key]["completed_var"].set(all_filled)
        self.sections_data[key]["completed"] = all_filled

    def _toggle_bilan2(self, key):
        enabled = bool(self.sections_widgets[key]["bilan2_var"].get())
        self.sections_data[key]["bilan2_enabled"] = enabled
        if enabled:
            self.sections_widgets[key]["bilan2_btn"].state(["!disabled"])
        else:
            self.sections_widgets[key]["bilan2_btn"].state(["disabled"])

    def _add_bilan(self, key, which=1):
        initial = self.sections_data[key]["bilan1" if which == 1 else "bilan2"]
        title = f"Saisir le {'1er' if which == 1 else '2e'} bilan - {SECTION_LABELS[key]}"
        text = self._prompt_multiline(title, initial)
        if text is not None:
            if which == 1:
                self.sections_data[key]["bilan1"] = text.strip()
            else:
                self.sections_data[key]["bilan2"] = text.strip()

    # ---- Couverture (aperçu mini) ----

    def _import_photo(self):
        p = filedialog.askopenfilename(
            title="Importer photo de l'élève",
            filetypes=[("Images", "*.png;*.jpg;*.jpeg;*.bmp")]
        )
        if p:
            self.photo_path = p
            self.update_cover_preview()

    def _import_section_photo(self, key):
        p = filedialog.askopenfilename(
            title=f"Photo pour {SECTION_LABELS[key]}",
            filetypes=[("Images", "*.png;*.jpg;*.jpeg;*.bmp")]
        )
        if p:
            self.sections_data[key]["photo"] = p
            self.sections_widgets[key]["photo_label"].configure(text=os.path.basename(p))
            self.update_cover_preview()

    def _mark_personal_completed(self):
        self.personal_completed = bool(
            self.nom_var.get().strip() and self.prenom_var.get().strip() and self.naissance_var.get().strip()
        )
        self.update_cover_preview()

    def _cover_canvas_size(self):
        # taille réelle du canvas de couverture
        cw = self.cover_canvas.winfo_width()
        ch = self.cover_canvas.winfo_height()
        if cw <= 1:
            cw = PREVIEW_WIDTH
        if ch <= 1:
            ch = 260
        return cw, ch

    def update_cover_preview(self):
        c = self.cover_canvas
        c.delete("all")
        cw, ch = self._cover_canvas_size()

        # Tente d'afficher la bannière top si disponible
        top_img_path = self._find_image_variant(os.path.join("img", "banniere-top.png"))
        if top_img_path and os.path.exists(top_img_path):
            try:
                pil = Image.open(top_img_path)
                ratio = pil.width / pil.height if pil.height else 1.0
                new_w = cw
                new_h = int(new_w / ratio)
                if new_h > min(160, int(ch * 0.5)):
                    new_h = min(160, int(ch * 0.5))
                    new_w = int(new_h * ratio)
                pil = pil.resize((new_w, new_h), Image.LANCZOS)
                tkimg = ImageTk.PhotoImage(pil)
                c.top_banner_image = tkimg
                c.create_image(0, 0, anchor="nw", image=tkimg)
                banner_h = new_h
            except Exception:
                banner_h = COVER_HEADER_HEIGHT
                c.create_rectangle(0, 0, cw, banner_h, fill=COVER_HEADER_COLOR, outline=COVER_HEADER_COLOR)
        else:
            banner_h = COVER_HEADER_HEIGHT
            c.create_rectangle(0, 0, cw, banner_h, fill=COVER_HEADER_COLOR, outline=COVER_HEADER_COLOR)
            c.create_text(cw // 2, banner_h // 2, text="PROFIL DE L'ELEVE", fill="white",
                          font=("Arial", 16, "bold"))

        # Zone texte à gauche (infos personnelles) avec fond olive
        x = TEXT_MARGIN_X
        y = banner_h + 10
        bg_w = min(480, int(cw * 0.55))
        bg_h = int(ch * 0.5)
        c.create_rectangle(x - 10, y - 8, x - 10 + bg_w, y - 8 + bg_h,
                           fill=COVER_PERSONAL_BG_PREVIEW, outline=COVER_PERSONAL_BG_PREVIEW)

        lines = []
        if self.nom_var.get().strip():
            lines.append(f"Nom: {self.nom_var.get().strip()}")
        if self.prenom_var.get().strip():
            lines.append(f"Prénom: {self.prenom_var.get().strip()}")
        if self.naissance_var.get().strip():
            lines.append(f"Date de naissance: {self.naissance_var.get().strip()}")

        ty = y
        for line in lines:
            c.create_text(x, ty, anchor="nw", text=line, font=("Arial", 12, "bold"), fill="white")
            ty += 24

        # Photo (si fournie) - mini-aperçu à droite
        if self.photo_path and os.path.exists(self.photo_path):
            try:
                max_side = min(160, int(ch * 0.55))
                pil = Image.open(self.photo_path)
                pil.thumbnail((max_side, max_side), Image.LANCZOS)
                tkimg = ImageTk.PhotoImage(pil)
                c.image = tkimg  # éviter le GC
                c.create_image(cw - max_side - 20, banner_h + 8, anchor="nw", image=tkimg)
            except Exception:
                pass

    # ---- Pagination & Aperçu ----

    def rebuild_pages_and_refresh(self):
        # Regroupe par domaine/sous-domaine et découpe en pages (pour l'aperçu)
        self.domain_page_map.clear()
        self.item_page_index.clear()
        # ordre des domaines fixe
        by_domain = OrderedDict((d, OrderedDict()) for d in self.domain_order)
        for it in self.selected_items:
            d = it.domain
            sd = it.subdomain or d
            by_domain.setdefault(d, OrderedDict())
            by_domain[d].setdefault(sd, [])
            by_domain[d][sd].append(it)

        for d in self.domain_order:
            pages = []
            current_page = []
            line_count = 0
            submap = by_domain.get(d, {})

            for sd, items in submap.items():
                # header sd
                if line_count + 1 > MAX_LINES_PER_SLIDE and current_page:
                    pages.append(current_page)
                    current_page = []
                    line_count = 0
                current_page.append((True, sd, None))
                line_count += 1

                # items
                for it in items:
                    if line_count + 1 > MAX_LINES_PER_SLIDE and current_page:
                        pages.append(current_page)
                        current_page = []
                        line_count = 0
                        # Répéter le header du sd sur la nouvelle page
                        current_page.append((True, sd, None))
                        line_count += 1
                    current_page.append((False, sd, it))
                    line_count += 1

            if current_page:
                pages.append(current_page)

            self.domain_page_map[d] = pages

            # indexer items -> page
            for pi, page in enumerate(pages):
                for is_header, sd, payload in page:
                    if not is_header and payload is not None:
                        self.item_page_index[payload.key()] = (d, pi)

        # Construire la liste plate de navigation (tous domaines)
        self.flat_pages = []
        for d in self.domain_order:
            pages = self.domain_page_map.get(d, [])
            if not pages:
                continue
            for pi in range(len(pages)):
                self.flat_pages.append((d, pi))

        # Ajuster current_flat_index
        if not self.flat_pages:
            self.current_flat_index = 0
            self.current_domain = None
        else:
            if self.current_flat_index >= len(self.flat_pages):
                self.current_flat_index = len(self.flat_pages) - 1
            self.current_domain = self.flat_pages[self.current_flat_index][0]

        self.update_preview()

    def prev_page(self):
        if not self.flat_pages:
            return
        if self.current_flat_index > 0:
            self.current_flat_index -= 1
            self.current_domain = self.flat_pages[self.current_flat_index][0]
            self.update_preview()

    def next_page(self):
        if not self.flat_pages:
            return
        if self.current_flat_index < len(self.flat_pages) - 1:
            self.current_flat_index += 1
            self.current_domain = self.flat_pages[self.current_flat_index][0]
            self.update_preview()

    def _preview_canvas_size(self):
        cw = self.preview_canvas.winfo_width()
        ch = self.preview_canvas.winfo_height()
        if cw <= 1:
            cw = PREVIEW_WIDTH
        if ch <= 1:
            ch = PREVIEW_HEIGHT
        return cw, ch

    def update_preview(self):
        # Met à jour la mini-couverture + la page domaine courante
        self.update_cover_preview()
        c = self.preview_canvas
        c.delete("all")

        cw, ch = self._preview_canvas_size()

        if not self.flat_pages:
            self.page_var.set("Page 0/0 — Aucune page (ajoutez des compétences)")
            c.create_text(cw//2, ch//2, text="Aucune page à afficher",
                          font=("Arial", 14, "italic"), fill="#666")
            return

        d, pi = self.flat_pages[self.current_flat_index]
        pages = self.domain_page_map.get(d, [])
        total_pages = len(self.flat_pages)
        self.page_var.set(f"Page {self.current_flat_index + 1}/{total_pages} — Domaine: {d} — p.{pi + 1}/{len(pages)}")

        ds = self.domain_states[d]
        # Bandeau de domaine
        c.create_rectangle(0, 0, cw, HEADER_HEIGHT, fill=ds.color, outline=ds.color)
        c.create_text(TEXT_MARGIN_X, HEADER_HEIGHT // 2, anchor="w",
                      text=d, fill=DEFAULT_TITLE_FG, font=("Arial", 16, "bold"))

        # Contenu
        y = HEADER_HEIGHT + 14
        x = TEXT_MARGIN_X
        body_font = ("Arial", ds.font_body[1])
        max_text_width = cw - 2 * TEXT_MARGIN_X - 10

        page = pages[pi] if 0 <= pi < len(pages) else []
        for is_header, sub, payload in page:
            if is_header:
                c.create_text(x, y, anchor="nw", text=sub, fill=ds.color,
                              font=("Arial", 13, "bold", "underline"))
                y += 22
            else:
                wrapped = self.wrap_text(payload.text, max_text_width, body_font)
                for li, line in enumerate(wrapped):
                    c.create_text(x + 16, y, anchor="nw",
                                  text=("• " + line if li == 0 else "  " + line),
                                  fill="black", font=body_font)
                    y += (ds.font_body[1] + LINE_SPACING)
                y += SUBHEADER_SPACING

        # Images de la page courante
        key = (d, pi)
        for img in self.page_images.get(key, []):
            c.create_image(img["pos"][0], img["pos"][1], image=img["tk"], anchor="nw")

    def wrap_text(self, text, max_width_px, font_tuple):
        f = tkfont.Font(family=font_tuple[0], size=font_tuple[1])
        words = text.split()
        lines = []
        cur = ""
        for w in words:
            test = w if not cur else cur + " " + w
            if f.measure(test) <= max_width_px:
                cur = test
            else:
                if cur:
                    lines.append(cur)
                cur = w
        if cur:
            lines.append(cur)
        return lines

    # ---- Sauvegarde / Chargement ----

    def save_project(self):
        path = filedialog.asksaveasfilename(title="Sauvegarder projet",
                                            defaultextension=".json",
                                            filetypes=[("JSON", "*.json")])
        if not path:
            return
        data = {
            "available": self.available,
            "domain_order": self.domain_order,
            "selected": [
                (it.domain, it.subdomain, it.text, it.ts, it.batch_id) for it in self.selected_items
            ],
            "domains": {
                d: {
                    "color": self.domain_states[d].color,
                    "font_body": self.domain_states[d].font_body,
                } for d in self.domain_order
            },
            "page_images": [
                {
                    "domain": d,
                    "page_index": pi,
                    "images": [
                        {"path": im["path"], "pos": im["pos"], "size": im["size"]}
                        for im in imgs
                    ]
                }
                for (d, pi), imgs in self.page_images.items()
            ],
            "infos": {
                "nom": self.nom_var.get(),
                "prenom": self.prenom_var.get(),
                "naissance": self.naissance_var.get(),
                "photo": self.photo_path,
                "personal_completed": self.personal_completed,
                "month": self.month_var.get(),
                "year": self.year_var.get(),
            },
            "sections": {
                key: {
                    "completed": self.sections_data[key]["completed"],
                    "fields": self.sections_data[key]["fields"],
                    "photo": self.sections_data[key]["photo"],
                    "bilan1": self.sections_data[key]["bilan1"],
                    "bilan2": self.sections_data[key]["bilan2"],
                    "bilan2_enabled": self.sections_data[key]["bilan2_enabled"],
                } for key in SECTION_KEYS
            }
        }
        try:
            with open(path, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
            messagebox.showinfo("Sauvegarde", f"Projet sauvegardé : {path}")
        except Exception as e:
            messagebox.showerror("Sauvegarde", str(e))

    def load_project(self):
        path = filedialog.askopenfilename(title="Charger projet", filetypes=[("JSON", "*.json")])
        if not path:
            return
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception as e:
            messagebox.showerror("Chargement", str(e))
            return

        # Restaurer domaines/compétences
        self.available = OrderedDict()
        for d, submap in data.get("available", {}).items():
            self.available[d] = OrderedDict()
            for sd, lst in submap.items():
                self.available[d][sd] = list(lst)

        self.domain_order = data.get("domain_order", [])
        self.domain_states.clear()
        domdata = data.get("domains", {})
        for idx, d in enumerate(self.domain_order):
            color = domdata.get(d, {}).get("color", DOMAIN_COLORS[idx % len(DOMAIN_COLORS)])
            ds = DomainState(d, color)
            ds.font_body = tuple(domdata.get(d, {}).get("font_body", DEFAULT_BODY_FONT))
            self.domain_states[d] = ds

        self.selected_items = []
        self.added_set = set()
        for tup in data.get("selected", []):
            # rétrocompat: (d, sd, txt) ou (d, sd, txt, ts, batch)
            if len(tup) == 3:
                d, sd, txt = tup
                ts = None
                batch_id = None
            else:
                d, sd, txt, ts, batch_id = tup
            it = CompetenceItem(d, sd, txt, ts=ts, batch_id=batch_id)
            self.selected_items.append(it)
            self.added_set.add(it.key())

        # Infos
        infos = data.get("infos", {})
        self.nom_var.set(infos.get("nom", ""))
        self.prenom_var.set(infos.get("prenom", ""))
        self.naissance_var.set(infos.get("naissance", ""))
        self.photo_path = infos.get("photo", None)
        self.personal_completed = bool(infos.get("personal_completed", False))
        self.month_var.set(infos.get("month", ""))
        self.year_var.set(infos.get("year", ""))

        # Sections
        secdata = data.get("sections", {})
        for key in SECTION_KEYS:
            sd = secdata.get(key, {})
            comp = bool(sd.get("completed", False))
            self.sections_data[key]["completed"] = comp
            if key in self.sections_widgets:
                self.sections_widgets[key]["completed_var"].set(comp)
            # fields
            fields = sd.get("fields", {})
            for fname in SECTION_FIELDS.keys():
                val = fields.get(fname, "")
                self.sections_data[key]["fields"][fname] = val
                if key in self.sections_widgets:
                    ent, var = self.sections_widgets[key]["entries"][fname]
                    var.set(val)
            # photo
            p = sd.get("photo", None)
            self.sections_data[key]["photo"] = p
            if key in self.sections_widgets:
                self.sections_widgets[key]["photo_label"].configure(text=os.path.basename(p) if p else "Aucune photo")
            # bilans
            self.sections_data[key]["bilan1"] = sd.get("bilan1", "")
            self.sections_data[key]["bilan2"] = sd.get("bilan2", "")
            self.sections_data[key]["bilan2_enabled"] = bool(sd.get("bilan2_enabled", False))
            if key in self.sections_widgets:
                self.sections_widgets[key]["bilan2_var"].set(self.sections_data[key]["bilan2_enabled"])
                if self.sections_data[key]["bilan2_enabled"]:
                    self.sections_widgets[key]["bilan2_btn"].state(["!disabled"])
                else:
                    self.sections_widgets[key]["bilan2_btn"].state(["disabled"])

        # Rebuild pages first to know page indices
        self.build_available_tree()
        self.refresh_selected_tree()
        self.rebuild_pages_and_refresh()

        # Restaurer images par page
        self.page_images.clear()
        for rec in data.get("page_images", []):
            d = rec.get("domain")
            pi = int(rec.get("page_index", 0))
            imgs = []
            for im in rec.get("images", []):
                p = im.get("path")
                pos = im.get("pos", [60, HEADER_HEIGHT + 30])
                size = im.get("size", [120, 120])
                try:
                    pil = Image.open(p).resize((int(size[0]), int(size[1])), Image.LANCZOS)
                    tkimg = ImageTk.PhotoImage(pil)
                    imgs.append({
                        "path": p, "pil": pil, "tk": tkimg,
                        "pos": [int(pos[0]), int(pos[1])],
                        "size": [int(size[0]), int(size[1])]
                    })
                except Exception:
                    pass
            if imgs:
                self.page_images[(d, pi)] = imgs

        self.update_preview()
        messagebox.showinfo("Chargement", "Projet chargé avec succès")

    # ---- Export PowerPoint ----

    def export_ppt(self):
        # Tente de (re)charger DOMAINES.txt pour descriptions
        self._load_domaines_descriptions()

        # Propose PRENOM_NOM.pptx comme nom initial
        suggested = self._default_ppt_filename()
        path = filedialog.asksaveasfilename(
            title="Créer PowerPoint",
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx")],
            initialfile=suggested
        )
        if not path:
            return
        try:
            prs = Presentation()
            # Page 1: Couverture
            self._is_cover_export = True
            self.build_cover_slide(prs)
            self._is_cover_export = False

            # Pages domaines
            self.rebuild_pages_and_refresh()

            for d in self.domain_order:
                pages = self.domain_page_map.get(d, [])
                if not pages:
                    continue
                ds = self.domain_states[d]
                for pi, page in enumerate(pages):
                    j = 0
                    current_sd = None
                    first_slide_for_page = True

                    # paramètres d'échelle basés sur la taille réelle de l'aperçu
                    cw, ch = self._preview_canvas_size()
                    preview_y_start = HEADER_HEIGHT + 14
                    preview_y_bottom_margin = 20
                    content_preview_height_px = max(1, ch - preview_y_start - preview_y_bottom_margin)

                    while j < len(page):
                        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blanc

                        # Description du domaine (si disponible)
                        domain_desc = (self.domain_descriptions.get(d) or "").strip()

                        # Bandeau domaine dynamique
                        banner_h = self.add_domain_banner(slide, prs, d, ds.color, domain_desc)

                        # Zone de contenu
                        left = Inches(0.6)
                        content_top = banner_h + Inches(0.1)
                        width = prs.slide_width - Inches(1.2)
                        height = prs.slide_height - content_top - Inches(0.9)

                        body_font_size = ds.font_body[1]
                        max_text_width_px = cw - 2 * TEXT_MARGIN_X - 10
                        y_px = preview_y_start
                        max_y_px = preview_y_start + content_preview_height_px
                        last_ts_slide = None

                        # Si on continue un sous-domaine sur une nouvelle diapo, réafficher son en-tête
                        if current_sd:
                            needed = 22
                            if y_px + needed <= max_y_px:
                                tb = slide.shapes.add_textbox(left, content_top + (y_px - preview_y_start) / content_preview_height_px * height, width, Inches(0.4))
                                tf = tb.text_frame
                                tf.clear()
                                p = tf.paragraphs[0]
                                p.text = current_sd
                                p.font.size = Pt(DEFAULT_BODY_SIZE_PT + 1)
                                p.font.bold = DEFAULT_SUBHEADER_BOLD
                                p.font.underline = DEFAULT_SUBHEADER_UNDERLINE
                                r, g, b = self.hex_to_rgb(self.domain_states[d].color)
                                p.font.color.rgb = RGBColor(r, g, b)
                                y_px += 22

                        while j < len(page):
                            is_header, sub, payload = page[j]
                            if is_header:
                                needed = 22
                                if y_px + needed > max_y_px:
                                    # Nouvelle diapo, on reprendra ce header
                                    current_sd = sub
                                    break
                                # Titre sous-domaine
                                tb = slide.shapes.add_textbox(left, content_top + (y_px - preview_y_start) / content_preview_height_px * height, width, Inches(0.4))
                                tf = tb.text_frame
                                tf.clear()
                                p = tf.paragraphs[0]
                                p.text = sub
                                p.font.size = Pt(DEFAULT_BODY_SIZE_PT + 1)
                                p.font.bold = DEFAULT_SUBHEADER_BOLD
                                p.font.underline = DEFAULT_SUBHEADER_UNDERLINE
                                r, g, b = self.hex_to_rgb(self.domain_states[d].color)
                                p.font.color.rgb = RGBColor(r, g, b)
                                y_px += 22
                                current_sd = sub
                                j += 1
                            else:
                                # Calcul de la hauteur nécessaire
                                ts = (payload.ts or "").strip()
                                prenom = (self.prenom_var.get() or "").strip()
                                full_text = f"• {prenom} {payload.text}".strip()
                                wrapped_lines = self.wrap_text(full_text, max_text_width_px, ("Arial", body_font_size))
                                lines_h_px = len(wrapped_lines) * (body_font_size + LINE_SPACING)
                                banner_h_px = 20 if (ts and ts != last_ts_slide) else 0
                                needed = banner_h_px + lines_h_px + SUBHEADER_SPACING

                                if y_px + needed > max_y_px:
                                    # Nouvelle diapo pour ce même élément
                                    break

                                # Bandeau date si nécessaire
                                if banner_h_px:
                                    band_shape = slide.shapes.add_shape(
                                        MSO_SHAPE.RECTANGLE,
                                        left,
                                        content_top + (y_px - preview_y_start) / content_preview_height_px * height,
                                        width, Inches(0.28)
                                    )
                                    band_shape.fill.solid()
                                    band_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
                                    try:
                                        band_shape.line.fill.background()
                                    except Exception:
                                        pass
                                    tb = slide.shapes.add_textbox(left,
                                                                  content_top + (y_px - preview_y_start) / content_preview_height_px * height,
                                                                  width, Inches(0.28))
                                    tf = tb.text_frame
                                    tf.clear()
                                    p = tf.paragraphs[0]
                                    p.text = ts
                                    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                                    p.font.size = Pt(10)
                                    p.font.bold = False
                                    p.font.color.rgb = RGBColor(255, 255, 255)
                                    y_px += banner_h_px
                                    last_ts_slide = ts

                                # Texte de la compétence
                                bullet_h = max(Inches(0.3), (lines_h_px / content_preview_height_px) * height)
                                tb = slide.shapes.add_textbox(left + Inches(0.2),
                                                              content_top + (y_px - preview_y_start) / content_preview_height_px * height,
                                                              width - Inches(0.2), bullet_h)
                                tf = tb.text_frame
                                tf.clear()
                                first_line = True
                                for li, line in enumerate(wrapped_lines):
                                    p = tf.paragraphs[0] if first_line else tf.add_paragraph()
                                    p.text = line
                                    p.font.size = Pt(body_font_size)
                                    p.font.bold = False
                                    p.font.color.rgb = RGBColor(0, 0, 0)
                                    if li == 0:
                                        p.level = 0
                                    first_line = False
                                y_px += lines_h_px
                                y_px += SUBHEADER_SPACING
                                j += 1

                        # Exporter images de la première diapo de cette page d'aperçu uniquement
                        if first_slide_for_page:
                            self.export_page_images(slide, prs, d, pi)
                            first_slide_for_page = False

                    # Diapos "Synthèse" par SECTION complétée
            for key in SECTION_KEYS:
                if self.sections_data[key]["completed"]:
                    self.build_section_synthesis_slide(prs, key)

            prs.save(path)
            messagebox.showinfo("Succès", f"PowerPoint sauvegardé : {path}")
        except Exception as e:
            messagebox.showerror("Export", f"Échec de l'export : {e}")

    def build_cover_slide(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blanc

        sw = prs.slide_width
        sh = prs.slide_height
        margin = Inches(0.6)

        # Bannière top (image si dispo)
        top_path = self._find_image_variant(os.path.join("img", "banniere-top.png"))
        used_banner_h = Inches(0.8)
        if top_path and os.path.exists(top_path):
            try:
                pic = slide.shapes.add_picture(top_path, Inches(0), Inches(0), width=sw)
                used_banner_h = pic.height
            except Exception:
                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), sw, Inches(0.8))
                rect.fill.solid()
                r, g, b = self.hex_to_rgb(COVER_HEADER_COLOR)
                rect.fill.fore_color.rgb = RGBColor(r, g, b)
                rect.line.fill.background()
                tb = slide.shapes.add_textbox(Inches(0), Inches(0), sw, Inches(0.8))
                tf = tb.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "PROFIL DE L'ELEVE"
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                used_banner_h = Inches(0.8)
        else:
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), sw, Inches(0.8))
            rect.fill.solid()
            r, g, b = self.hex_to_rgb(COVER_HEADER_COLOR)
            rect.fill.fore_color.rgb = RGBColor(r, g, b)
            rect.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0), Inches(0), sw, Inches(0.8))
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "PROFIL DE L'ELEVE"
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            used_banner_h = Inches(0.8)

        content_top = used_banner_h + Inches(0.15)

        # Colonne principale: fond olive + infos personnelles
        left_left = margin
        left_w = sw - 2 * margin
        left_top = content_top
        left_h = Inches(1.8)

        rect_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_left, left_top, left_w, left_h)
        rect_bg.fill.solid()
        rect_bg.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_3
        rect_bg.line.fill.background()

        tb2 = slide.shapes.add_textbox(left_left + Inches(0.15), left_top + Inches(0.12), left_w - Inches(0.3), left_h - Inches(0.24))
        tf2 = tb2.text_frame
        tf2.clear()
        lines = []
        if self.nom_var.get().strip():
            lines.append(f"Nom: {self.nom_var.get().strip()}")
        if self.prenom_var.get().strip():
            lines.append(f"Prénom: {self.prenom_var.get().strip()}")
        if self.naissance_var.get().strip():
            lines.append(f"Date de naissance: {self.naissance_var.get().strip()}")

        first = True
        for line in lines:
            if first:
                p = tf2.paragraphs[0]
                p.text = line
                first = False
            else:
                p = tf2.add_paragraph()
                p.text = line
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

        # Photo élève (en haut à droite)
        if self.photo_path and os.path.exists(self.photo_path):
            try:
                max_photo_h = Inches(1.4)
                pic = slide.shapes.add_picture(self.photo_path, Inches(0), Inches(0), height=max_photo_h)
                pic.left = left_left + left_w - pic.width - Inches(0.2)
                pic.top = left_top + Inches(0.2)
            except Exception:
                pass

        # Présentation HORIZONTALE des sections
        row_top = left_top + left_h + Inches(0.2)
        col_count = len(SECTION_KEYS)
        if col_count < 1:
            self.add_bottom_banner(slide, prs)
            return
        col_w = (sw - 2 * margin) / col_count
        col_text_h = Inches(1.0)
        col_photo_h = Inches(1.6)

        for idx, key in enumerate(SECTION_KEYS):
            col_left = margin + col_w * idx
            section_title = SECTION_LABELS[key]
            fields = self.sections_data[key]["fields"]

            # Bloc titre + 3 lignes
            tb = slide.shapes.add_textbox(col_left, row_top, col_w, col_text_h)
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = section_title
            p.font.bold = True
            p.font.size = Pt(12)

            for fname, flabel in SECTION_FIELDS.items():
                val = fields.get(fname, "").strip()
                if not val:
                    continue
                sp = tf.add_paragraph()
                sp.text = f"{flabel}: {val}"
                sp.level = 1
                sp.font.size = Pt(10)

            # Emplacement photo sous le bloc texte
            ph_top = row_top + col_text_h + Inches(0.05)
            ph_w = col_w
            ph_h = col_photo_h

            ph_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_left, ph_top, ph_w, ph_h)
            ph_rect.fill.solid()
            ph_rect.fill.fore_color.rgb = RGBColor(245, 245, 245)
            ph_rect.line.fill.background()

            sec_photo = self.sections_data[key]["photo"]
            if sec_photo and os.path.exists(sec_photo):
                try:
                    with Image.open(sec_photo) as im:
                        iw, ih = im.size
                        box_w = ph_w
                        box_h = ph_h
                        img_ratio = iw / ih if ih else 1.0
                        box_ratio = box_w / box_h if box_h else 1.0
                        if img_ratio >= box_ratio:
                            pic = slide.shapes.add_picture(sec_photo, col_left, ph_top, width=box_w)
                            pic.top = ph_top + (box_h - pic.height) // 2
                        else:
                            pic = slide.shapes.add_picture(sec_photo, col_left, ph_top, height=box_h)
                            pic.left = col_left + (box_w - pic.width) // 2
                except Exception:
                    pass

        self.add_bottom_banner(slide, prs)

    def add_domain_banner(self, slide, prs, domain_name, color_hex, domain_desc=""):
        """
        Crée le bandeau supérieur de domaine (rectangle coloré + titre + description).
        Le bandeau s'AGRANDIT automatiquement pour que la description ne déborde pas.
        Retourne la hauteur du bandeau (Length).
        """
        sw = prs.slide_width

        # Paramètres de mise en page
        title_left = Inches(0.4)
        title_top = Inches(0.05)
        title_height = Inches(0.6)
        text_width = sw - Inches(0.8)
        desc_font_pt = 8

        # Estimation du nombre de lignes pour la description, en se basant sur l'aperçu
        def compute_desc_lines(desc_text: str) -> list[str]:
            if not desc_text.strip():
                return []
            # approx: largeur disponible proportionnelle au canvas de preview
            cw, _ = self._preview_canvas_size()
            max_width_px = max(50, int(cw) - 40)
            return self.wrap_text(desc_text, max_width_px, ("Arial", desc_font_pt))

        desc_lines = compute_desc_lines(domain_desc)
        # Hauteur des lignes en pouces: approx 1.25 x taille (en points/72)
        line_height_in = (desc_font_pt / 72.0) * 1.25
        desc_height = Inches(line_height_in * max(1, len(desc_lines))) if desc_lines else Inches(0)

        # Position par défaut de la description (sous le titre)
        desc_top = Inches(0.65)

        # Hauteur du bandeau minimum
        min_banner_h = Inches(1.2)
        # Hauteur nécessaire pour contenir la description (si présente)
        needed_banner_h = desc_top + desc_height + Inches(0.1)
        banner_h = max(min_banner_h, needed_banner_h)

        # Rectangle bandeau
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            sw, banner_h
        )
        rect.fill.solid()
        r, g, b = self.hex_to_rgb(color_hex)
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        rect.line.fill.background()

        # Titre du domaine
        tb = slide.shapes.add_textbox(
            title_left, title_top,
            text_width, title_height
        )
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = domain_name
        p.font.size = Pt(DEFAULT_TITLE_SIZE_PT)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Description
        if desc_lines:
            desc_tb = slide.shapes.add_textbox(title_left, desc_top, text_width, desc_height)
            desc_tf = desc_tb.text_frame
            desc_tf.clear()
            desc_tf.word_wrap = True
            first = True
            for line in desc_lines:
                if first:
                    pd = desc_tf.paragraphs[0]
                    first = False
                else:
                    pd = desc_tf.add_paragraph()
                pd.text = line
                pd.font.size = Pt(desc_font_pt)
                pd.font.bold = False
                pd.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                pd.font.color.rgb = RGBColor(255, 255, 255)

        return banner_h

    def export_page_images(self, slide, prs, domain, page_index):
        # Map coordonnées apercu -> slide pour la page (domain, page_index)
        key = (domain, page_index)
        sw = prs.slide_width
        sh = prs.slide_height

        cw, ch = self._preview_canvas_size()

        for img in self.page_images.get(key, []):
            lx = int(img["pos"][0] / cw * sw)
            ly = int(img["pos"][1] / ch * sh)
            w = int(img["size"][0] / cw * sw)
            h = int(img["size"][1] / ch * sh)
            try:
                if w > 0 and h > 0:
                    slide.shapes.add_picture(img["path"], lx, ly, width=w, height=h)
            except Exception:
                pass

    # ---- Diapo Synthèse par SECTION ----

    def build_section_synthesis_slide(self, prs, key):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blanc

        sw = prs.slide_width
        sh = prs.slide_height

        # Titre "Synthèse"
        band_h = Inches(0.8)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), sw, band_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 0, 0)  # rouge
        rect.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0), Inches(0), sw, band_h)
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Synthèse"
        p.font.size = Pt(20)
        p.font.bold = False
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # Sous-titre = intitulé de la SECTION
        subtitle = SECTION_LABELS.get(key, key)
        tb2 = slide.shapes.add_textbox(Inches(0.6), band_h + Inches(0.2), sw - Inches(1.2), Inches(0.6))
        tf2 = tb2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.bold = False
        p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # Zones de bilans
        top_content = band_h + Inches(1.0)
        left = Inches(0.6)
        width = sw - Inches(1.2)
        available_h = sh - top_content - Inches(0.9)

        bilan1 = (self.sections_data[key]["bilan1"] or "").strip()
        use_bilan2 = bool(self.sections_data[key]["bilan2_enabled"])
        bilan2 = (self.sections_data[key]["bilan2"] or "").strip() if use_bilan2 else ""

        if use_bilan2:
            box_h = available_h / 2.0 - Inches(0.2)
            # Bilan 1
            self._add_bilan_box(slide, left, top_content, width, box_h, "Bilan", bilan1)
            # Bilan 2
            self._add_bilan_box(slide, left, top_content + box_h + Inches(0.2), width, box_h, "Bilan", bilan2)
        else:
            self._add_bilan_box(slide, left, top_content, width, available_h, "Bilan", bilan1)

    def _add_bilan_box(self, slide, left, top, width, height, title, content):
        # Titre "Bilan"
        title_tb = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        title_tf = title_tb.text_frame
        title_tf.clear()
        pt = title_tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(14)
        pt.font.bold = True
        # Contenu
        text_tb = slide.shapes.add_textbox(left, top + Inches(0.45), width, max(Inches(0.8), height - Inches(0.45)))
        text_tf = text_tb.text_frame
        text_tf.clear()
        text_tf.word_wrap = True
        if content:
            lines = content.splitlines()
            first = True
            for line in lines:
                if first:
                    p = text_tf.paragraphs[0]
                    p.text = line
                    first = False
                else:
                    p = text_tf.add_paragraph()
                    p.text = line
                p.font.size = Pt(12)
        else:
            p = text_tf.paragraphs[0]
            p.text = ""
            p.font.size = Pt(12)

    # ---- Bannières top/bas utilitaires ----

    def add_bottom_banner(self, slide, prs):
        # Affichage uniquement sur la page de garde
        if not getattr(self, "_is_cover_export", False):
            return
        # Ajoute une image de bannière en bas si trouvée
        path_try = [
            os.path.join("img", "banniere-bas.png"),
            os.path.join("img", "banniere-bas.jpg"),
            os.path.join("img", "banniere-bas.jpeg"),
        ]
        img_path = next((p for p in path_try if os.path.exists(p)), None)
        if not img_path:
            return
        try:
            sw = prs.slide_width
            sh = prs.slide_height
            pic = slide.shapes.add_picture(img_path, Inches(0), sh - Inches(0.5), width=sw)
            pic.top = sh - pic.height
        except Exception:
            pass

    def _find_image_variant(self, path_with_default_ext):
        # Si le chemin donné existe, l'utiliser, sinon essayer variantes jpg/jpeg
        if os.path.exists(path_with_default_ext):
            return path_with_default_ext
        base, ext = os.path.splitext(path_with_default_ext)
        for e in [".png", ".jpg", ".jpeg"]:
            if os.path.exists(base + e):
                return base + e
        return None

    # ---- Lecture DOMAINES.txt ----

    def _load_domaines_descriptions(self):
        path = os.path.join(os.getcwd(), "DOMAINES.txt")
        if not os.path.exists(path):
            return
        try:
            domain = None
            sub = None
            buf = []
            self.domain_descriptions.clear()
            self.subdomain_descriptions.clear()

            def commit():
                nonlocal domain, sub, buf
                if not domain:
                    buf = []
                    return
                text = "\n".join(buf).strip()
                if not text:
                    buf = []
                    return
                if sub:
                    self.subdomain_descriptions[(domain, sub)] = text
                else:
                    self.domain_descriptions[domain] = text
                buf = []

            with open(path, "r", encoding="utf-8-sig") as f:
                for raw in f:
                    line = raw.rstrip("\n")
                    if not line.strip():
                        if buf is not None:
                            buf.append("")
                        continue
                    cleaned = line.replace("\u202f", " ").replace("\u00a0", " ").replace("\u2019", "'").strip()

                    if cleaned.startswith("##-"):
                        # Commit précédent
                        commit()
                        # Nouveau domaine
                        d = cleaned[3:].strip()
                        if d.lower().startswith("domaine"):
                            parts = d.split(None, 1)
                            d = parts[1] if len(parts) > 1 else d
                        domain = d
                        sub = None
                        buf = []
                    elif cleaned.startswith("#-"):
                        # Commit précédent (domaine ou sous-domaine précédent)
                        commit()
                        s = cleaned[2:].strip()
                        if s.lower().startswith("sous-domaine:"):
                            s = s.split(":", 1)[1].strip()
                        sub = s
                        buf = []
                    else:
                        buf.append(cleaned)
                commit()
        except Exception as e:
            messagebox.showwarning("DOMAINES.txt", f"Impossible de lire DOMAINES.txt : {e}")

    # ---- Utils ----

    @staticmethod
    def hex_to_rgb(hx):
        hx = hx.lstrip("#")
        return tuple(int(hx[i:i+2], 16) for i in (0, 2, 4))

    @staticmethod
    def sanitize_filename(name: str) -> str:
        invalid = '<>:"/\\|?*'
        for ch in invalid:
            name = name.replace(ch, "_")
        name = name.strip().replace(" ", "_")
        name = "".join(c if (c.isalnum() or c in ("_", "-")) else "_" for c in name)
        while "__" in name:
            name = name.replace("__", "_")
        return name.strip("_") or "presentation"

    def _default_ppt_filename(self) -> str:
        nom = (self.nom_var.get() or "").strip()
        prenom = (self.prenom_var.get() or "").strip()
        # PRENOM_NOM comme demandé
        base = f"{prenom}_{nom}".strip("_") if (nom or prenom) else "presentation"
        base = self.sanitize_filename(base)
        return f"{base}.pptx"

    # ---- Fenêtre de saisie multiligne ----
    def _prompt_multiline(self, title, initial_text=""):
        top = tk.Toplevel(self.root)
        top.title(title)
        top.transient(self.root)
        top.grab_set()
        top.geometry("800x600")
        top.minsize(700, 520)

        frm = ttk.Frame(top)
        frm.pack(fill="both", expand=True, padx=8, pady=8)

        txt = tk.Text(frm, wrap="word")
        txt.pack(fill="both", expand=True)
        if initial_text:
            txt.insert("1.0", initial_text)

        btns = ttk.Frame(frm)
        btns.pack(fill="x", pady=6)
        result = {"text": None}

        def on_ok():
            result["text"] = txt.get("1.0", "end-1c")
            top.destroy()

        def on_cancel():
            result["text"] = None
            top.destroy()

        ttk.Button(btns, text="OK", command=on_ok).pack(side="right", padx=4)
        ttk.Button(btns, text="Annuler", command=on_cancel).pack(side="right", padx=4)

        top.wait_window()
        return result["text"]


# ==== Lancement ====

if __name__ == "__main__":
    root = tk.Tk()
    app = CompetenceApp(root)
    root.mainloop()
